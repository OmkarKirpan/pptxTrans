import os
import uuid
import logging
import time
import asyncio
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime
from pptx import Presentation
from PIL import Image
from PIL import ImageDraw
from PIL import ImageFont
import json
import shutil
import xml.etree.ElementTree as ET
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Emu, Pt
import base64
from pptx.enum.shapes import MSO_SHAPE_TYPE
import subprocess
import tempfile
import glob

from app.models.schemas import (
    ProcessedSlide,
    SlideShape,
    ProcessedPresentation,
    OverallProcessingStatus,
    ShapeType,
    CoordinateUnit,
    ProcessingStatusResponse,
    ProcessingStatus
)
from app.services.supabase_service import upload_file_to_supabase, update_job_status
from app.services.job_status import update_job_status as update_local_job_status, get_job_status
# Import get_settings instead of settings
from app.core.config import get_settings

logger = logging.getLogger(__name__)

# Get settings
settings = get_settings()

# Removed hardcoded LIBREOFFICE_PATH, will use settings.LIBREOFFICE_PATH

# Dictionary to keep track of job file paths for retry capability
_job_file_paths = {}


async def get_job_file_path(job_id: str) -> Optional[str]:
    """
    Get the file path for a job.
    Returns None if the job doesn't exist or the file has been cleaned up.
    """
    # First check our in-memory cache
    if job_id in _job_file_paths and os.path.exists(_job_file_paths[job_id]):
        return _job_file_paths[job_id]

    # If not in memory, we can't retrieve it since we don't persistently store file paths
    return None


async def _generate_svgs_for_all_slides_libreoffice(
    presentation_path: str,
    output_dir: str,  # This will be the directory where final slide_N.svg files are stored
    slide_count: int
) -> Dict[int, str]:
    """
    Convert each slide from a presentation to an individual SVG file using LibreOffice.
    Iterates through slides and calls LibreOffice for each one.
    """
    if not settings.LIBREOFFICE_PATH or not os.path.exists(settings.LIBREOFFICE_PATH):
        logger.warning(
            "LibreOffice path not configured or invalid. Skipping LibreOffice SVG generation.")
        return {}

    abs_presentation_path = os.path.abspath(presentation_path)
    if not os.path.exists(abs_presentation_path):
        logger.error(
            f"Absolute presentation path not found: {abs_presentation_path}")
        return {}

    # Temporary directory for LibreOffice to write individual SVG outputs before renaming/moving
    # This main temp dir will contain uniquely named SVGs from each LO call.
    temp_lo_svg_output_dir = os.path.join(
        output_dir, f"lo_indiv_svg_temp_{uuid.uuid4().hex[:8]}")
    os.makedirs(temp_lo_svg_output_dir, exist_ok=True)
    abs_temp_lo_svg_output_dir = os.path.abspath(temp_lo_svg_output_dir)

    generated_svg_paths: Dict[int, str] = {}

    for i in range(slide_count):
        slide_number = i + 1  # Slide numbers are 1-based
        # LibreOffice page numbers for export are often 1-based.
        # Output filename from LO will usually be <original_filename_without_ext>_page_<page_number>.svg
        # or just <original_filename_without_ext>.svg if it can only output one file.
        # We will try to make the output file specific to avoid overwrites in the temp dir.

        # The actual output filename by LibreOffice when converting a single page from a multi-page doc
        # with --export-page and --outdir can be just the original filename with new extension.
        # To handle this safely and ensure unique files in our temp dir from each call,
        # we can tell LO to output to a sub-directory per slide, or rename immediately.
        # Simpler: use outdir, and expect LO to name the file as presentation_name.svg.
        # We then rename this to our slide_N.svg convention.

        # Let LO write its default name (e.g., original_filename.svg) into the temp output dir.
        # We will move and rename it after successful conversion of *this specific slide*.

        # Construct the expected output file name by LibreOffice (usually original_filename.svg)
        presentation_basename = os.path.splitext(
            os.path.basename(abs_presentation_path))[0]
        expected_lo_output_filename = f"{presentation_basename}.svg"
        # This is the path where LO will place its output for the current slide conversion
        current_slide_lo_output_path = os.path.join(
            abs_temp_lo_svg_output_dir, expected_lo_output_filename)

        # Delete this expected output file if it exists from a previous iteration (unlikely with unique temp dir per call now)
        # but good for safety if LO overwrites.
        if os.path.exists(current_slide_lo_output_path):
            try:
                os.remove(current_slide_lo_output_path)
            except OSError as e:
                logger.warning(
                    f"Could not remove existing temp LO output {current_slide_lo_output_path}: {e}")

        try:
            logger.info(
                f"Attempting to convert slide {slide_number} of {slide_count} from {abs_presentation_path} to SVG")

            # Command to export a single page (slide)
            # Using --export-filter-options="PageRange=<page_num>" might be more robust for impress
            # or simply --page <page_num> or --select-page <page_num> (syntax varies)
            # The most common for Impress seems to be an export filter option like PageRange=N
            # Let's try with a filter option. Page numbers are usually 1-based.
            # Filter options format: PageRange=N for a single page N (1-based).
            # Or PageRange=N-M for a range. We need PageRange=slide_number.
            # The filter name is impress_svg_Export. Options are appended after a colon.
            # Example: "impress_svg_Export:PageRange=1"
            # For SVG, some use "impress_svg_Export:SVGPages=1" (1 for current, 2 for all)
            # The --page option in soffice man page: --page <range>
            # Example: --page 1-1 for first page. Let's try this as it's simpler.

            command = [
                settings.LIBREOFFICE_PATH,
                "--headless",
                # "--convert-to", f'svg:impress_svg_Export:PageRange={slide_number}', # This is complex and might not be right
                "--convert-to", "svg:impress_svg_Export",  # Keep filter simple
                # Added export-page argument
                f"--export-page", str(slide_number),
                "--outdir", abs_temp_lo_svg_output_dir,  # Output to our general temp LO dir
                abs_presentation_path
            ]

            process = subprocess.run(
                command,
                check=True,
                capture_output=True,
                text=True,
                timeout=120  # Shorter timeout for single slide
            )
            logger.info(
                f"LibreOffice SVG conversion for slide {slide_number} stdout: {process.stdout}")
            if process.stderr:
                logger.warning(
                    f"LibreOffice SVG conversion for slide {slide_number} stderr: {process.stderr}")

            # After conversion, LibreOffice should have created a file (e.g., originalfilename.svg) in abs_temp_lo_svg_output_dir
            if os.path.exists(current_slide_lo_output_path):
                final_svg_name = f"slide_{slide_number}.svg"
                # Place in final output_dir
                final_svg_path = os.path.join(output_dir, final_svg_name)
                shutil.move(current_slide_lo_output_path, final_svg_path)
                generated_svg_paths[slide_number] = final_svg_path
                logger.info(
                    f"Successfully converted and moved slide {slide_number} to {final_svg_path}")
            else:
                logger.error(
                    f"LibreOffice converted slide {slide_number}, but expected output file {current_slide_lo_output_path} not found.")
                # If one slide fails, we might want to stop or continue and use fallback for this one.
                # For now, let's return empty to trigger fallback for all if any single slide fails this way.
                # return {} # This would stop all LO processing
                # Better: just skip this slide, allow fallback for it later
                continue

        except subprocess.CalledProcessError as e:
            logger.error(
                f"Error running LibreOffice for SVG conversion of slide {slide_number}: {e}")
            logger.error(
                f"SVG Command (slide {slide_number}) output: {e.stdout}")
            logger.error(
                f"SVG Command (slide {slide_number}) error: {e.stderr}")
            continue  # Continue to next slide, allow fallback for this one
        except subprocess.TimeoutExpired:
            logger.error(
                f"LibreOffice SVG conversion for slide {slide_number} timed out.")
            continue  # Continue to next slide, allow fallback for this one
        except Exception as e:
            logger.error(
                f"Unexpected error during LibreOffice SVG conversion for slide {slide_number}: {str(e)}", exc_info=True)
            continue  # Continue to next slide, allow fallback for this one

    # Clean up the main temporary directory for LO outputs if it still exists and is empty
    # (individual files should have been moved or handled)
    if os.path.exists(abs_temp_lo_svg_output_dir):
        try:
            if not os.listdir(abs_temp_lo_svg_output_dir):  # Check if empty
                shutil.rmtree(abs_temp_lo_svg_output_dir)
                logger.info(
                    f"Cleaned up empty temporary LibreOffice individual SVG output directory: {abs_temp_lo_svg_output_dir}")
            else:
                logger.warning(
                    f"Temporary LibreOffice individual SVG output directory {abs_temp_lo_svg_output_dir} is not empty. Manual check may be needed.")
        except Exception as e:
            logger.error(
                f"Error cleaning up temp LO individual SVG dir {abs_temp_lo_svg_output_dir}: {e}")

    return generated_svg_paths


async def queue_pptx_processing(
    job_id: str,
    session_id: str,
    file_path: str,
    source_language: Optional[str] = None,
    target_language: Optional[str] = None,
    generate_thumbnails: bool = True
) -> None:
    """
    Queue the PPTX processing task.
    """
    # Store the file path for potential retry
    _job_file_paths[job_id] = file_path

    await update_local_job_status(
        job_id=job_id,
        status=ProcessingStatusResponse(
            job_id=job_id,
            session_id=session_id,
            status=ProcessingStatus.PROCESSING,
            progress=0,
            current_stage="Starting processing"
        )
    )
    loop = asyncio.get_event_loop()
    loop.create_task(
        process_pptx(
            job_id=job_id,
            session_id=session_id,
            file_path=file_path,
            source_language=source_language,
            target_language=target_language,
            generate_thumbnails=generate_thumbnails
        )
    )


async def process_pptx(
    job_id: str,
    session_id: str,
    file_path: str,
    source_language: Optional[str] = None,
    target_language: Optional[str] = None,
    generate_thumbnails: bool = True
) -> None:
    """
    Process a PPTX file.
    """
    start_time = time.time()
    # processing_dir is the main directory for this job's outputs (SVGs, JSON)
    # It's inside the uploaded file's directory.
    uploaded_file_dir = os.path.dirname(file_path)
    processing_output_dir = os.path.join(
        uploaded_file_dir, "processing_output")
    os.makedirs(processing_output_dir, exist_ok=True)

    libreoffice_svgs: Dict[int, str] = {}

    try:
        # Check LibreOffice availability
        if settings.LIBREOFFICE_PATH and os.path.exists(settings.LIBREOFFICE_PATH):
            try:
                # Use --help as it's more reliable than --version for some LO installations
                test_command = [settings.LIBREOFFICE_PATH, "--help"]
                test_result = subprocess.run(
                    test_command, check=True, capture_output=True, text=True, timeout=30)
                logger.info(
                    f"LibreOffice is available at: {settings.LIBREOFFICE_PATH}")
            except Exception as e:
                logger.warning(
                    f"LibreOffice path is set but test failed: {str(e)}. Will use fallback SVG generation if optimized path fails.")
        else:
            logger.info(
                "LibreOffice path not configured or invalid. Using fallback SVG generation.")

        presentation = Presentation(file_path)
        slide_count = len(presentation.slides)

        await update_local_job_status(
            job_id=job_id,
            status=ProcessingStatusResponse(
                job_id=job_id, session_id=session_id, status=ProcessingStatus.PROCESSING,
                progress=5, current_stage=f"Opened presentation with {slide_count} slides"
            )
        )

        # Attempt to generate all SVGs using LibreOffice in one go
        if settings.LIBREOFFICE_PATH and os.path.exists(settings.LIBREOFFICE_PATH):
            # Pass the original file_path to LibreOffice, and processing_output_dir for its outputs
            libreoffice_svgs = await _generate_svgs_for_all_slides_libreoffice(
                file_path, processing_output_dir, slide_count
            )
            if libreoffice_svgs:
                logger.info(
                    f"Successfully pre-generated {len(libreoffice_svgs)} SVGs using LibreOffice.")
            else:
                logger.info(
                    "Failed to pre-generate SVGs with LibreOffice, will use fallback per slide.")

        processed_slides_data = []
        for idx, slide in enumerate(presentation.slides):
            slide_number = idx + 1
            progress = 5 + int((idx / slide_count) * 90)
            await update_local_job_status(
                job_id=job_id,
                status=ProcessingStatusResponse(
                    job_id=job_id, session_id=session_id, status=ProcessingStatus.PROCESSING,
                    progress=progress, current_stage=f"Processing slide {slide_number} of {slide_count}"
                )
            )

            # slide_specific_output_dir is for thumbnails and fallback SVGs for this specific slide
            slide_specific_output_dir = os.path.join(
                processing_output_dir, f"slide_{slide_number}_assets")
            os.makedirs(slide_specific_output_dir, exist_ok=True)

            processed_slide = await process_slide(
                slide=slide,
                slide_number=slide_number,
                # Pass slide_specific_output_dir for assets related to this slide (thumbnails, fallback SVG)
                slide_assets_dir=slide_specific_output_dir,
                # Pass processing_output_dir for the main LibreOffice SVGs if available
                main_processing_dir=processing_output_dir,
                libreoffice_generated_svg_path=libreoffice_svgs.get(
                    slide_number),  # Pass path if LO generated it
                session_id=session_id,
                generate_thumbnail=generate_thumbnails
            )
            processed_slides_data.append(processed_slide)

        processing_time = int(time.time() - start_time)
        result = ProcessedPresentation(
            session_id=session_id, slide_count=slide_count,
            processing_status=OverallProcessingStatus.COMPLETED,
            processing_time=processing_time, slides=processed_slides_data
        )

        result_file = os.path.join(
            processing_output_dir, f"result_{session_id}.json")
        with open(result_file, "w") as f:
            # Use result.dict() for Pydantic v1, result.model_dump() for v2
            json.dump(result.dict(), f, indent=4)

        result_url = await upload_file_to_supabase(
            file_path=result_file,
            bucket="processing-results", destination_path=f"{session_id}/result.json"
        )

        await update_local_job_status(
            job_id=job_id,
            status=ProcessingStatusResponse(
                job_id=job_id, session_id=session_id, status=ProcessingStatus.COMPLETED,
                progress=100, current_stage="Processing completed", completed_at=datetime.now()
            )
        )
        await update_job_status(
            session_id=session_id, status="completed",
            slide_count=slide_count, result_url=result_url
        )

    except Exception as e:
        logger.error(f"Error processing PPTX: {str(e)}", exc_info=True)
        await update_local_job_status(
            job_id=job_id,
            status=ProcessingStatusResponse(
                job_id=job_id, session_id=session_id, status=ProcessingStatus.FAILED,
                progress=0, current_stage="Processing failed", error=str(e)
            )
        )
        await update_job_status(
            session_id=session_id, status="failed", error=str(e)
        )
    finally:
        # Clean up the main directory containing the uploaded file and its processing_output
        try:
            # uploaded_file_dir is the parent of processing_output_dir and contains the original upload
            # This was: shutil.rmtree(os.path.dirname(file_path)) which is uploaded_file_dir
            if os.path.exists(uploaded_file_dir):
                # Only clean up if this is not a retry attempt (otherwise we'd lose the file)
                status = await get_job_status(job_id)
                if status and status.status != "queued" and status.status != "processing":
                    shutil.rmtree(uploaded_file_dir)
                    logger.info(
                        f"Cleaned up temporary processing directory: {uploaded_file_dir}")

                    # Remove from job file paths cache if we've deleted the file
                    if job_id in _job_file_paths:
                        del _job_file_paths[job_id]
        except Exception as e:
            logger.error(
                f"Error cleaning up temporary files at {uploaded_file_dir}: {str(e)}")


async def process_slide(
    slide,  # This is a python-pptx Slide object
    slide_number: int,
    slide_assets_dir: str,  # Dir for thumbnails, fallback SVGs for this specific slide
    # Main dir where LO SVGs might be (e.g. slide_1.svg)
    main_processing_dir: str,
    # Path to LO SVG if pre-generated
    libreoffice_generated_svg_path: Optional[str],
    session_id: str,
    generate_thumbnail: bool = True
) -> ProcessedSlide:
    """
    Process a single slide.
    Uses pre-generated LibreOffice SVG if available, otherwise falls back to ElementTree.
    """
    slide_id = str(uuid.uuid4())
    # slide_assets_dir is already created by process_pptx
    # os.makedirs(slide_assets_dir, exist_ok=True)

    # Get slide dimensions from the presentation's slide master
    # In python-pptx, slides inherit dimensions from slide masters
    presentation = slide.part.package.presentation_part.presentation
    slide_width_emu = presentation.slide_width
    slide_height_emu = presentation.slide_height

    # Extract shapes and their data first, as it's needed for both SVG fallback and final output
    extracted_shapes_data = extract_shapes(
        slide, slide_width_emu, slide_height_emu)

    svg_file_to_upload: Optional[str] = None

    if libreoffice_generated_svg_path and os.path.exists(libreoffice_generated_svg_path):
        logger.info(
            f"Using pre-generated LibreOffice SVG for slide {slide_number}: {libreoffice_generated_svg_path}")
        svg_file_to_upload = libreoffice_generated_svg_path
    else:
        logger.info(
            f"LibreOffice SVG not available for slide {slide_number}. Generating SVG using ElementTree fallback.")
        fallback_svg_path = os.path.join(
            slide_assets_dir, f"slide_{slide_number}_fallback.svg")
        # Pass extracted_shapes_data to avoid re-calculating
        create_svg_from_slide(
            slide_shapes_data=extracted_shapes_data,
            file_path=fallback_svg_path,
            width_emu=slide_width_emu,
            height_emu=slide_height_emu,
            slide_background_fill=get_slide_background_fill(
                slide)  # Get background fill
        )
        if os.path.exists(fallback_svg_path):
            svg_file_to_upload = fallback_svg_path
        else:
            logger.error(
                f"Fallback SVG generation failed for slide {slide_number}")
            # Create a minimal empty SVG as a last resort to avoid crashes downstream
            fallback_svg_path = os.path.join(
                slide_assets_dir, f"slide_{slide_number}_empty.svg")
            create_minimal_svg(fallback_svg_path,
                               slide_width_emu, slide_height_emu)
            svg_file_to_upload = fallback_svg_path

    svg_url = None
    if svg_file_to_upload and os.path.exists(svg_file_to_upload):
        svg_url = await upload_file_to_supabase(
            file_path=svg_file_to_upload,
            bucket="slide-visuals", destination_path=f"{session_id}/slide_{slide_number}.svg"
        )
    else:
        logger.error(
            f"No SVG file could be prepared for upload for slide {slide_number}")

    thumbnail_url = None
    if generate_thumbnail:
        thumbnail_file = os.path.join(
            slide_assets_dir, f"thumbnail_{slide_number}.png")
        # Pass shapes_data if create_thumbnail_from_slide can use it, or slide object
        create_thumbnail_from_slide_pil(
            slide, extracted_shapes_data, thumbnail_file, slide_width_emu, slide_height_emu)
        if os.path.exists(thumbnail_file):
            thumbnail_url = await upload_file_to_supabase(
                file_path=thumbnail_file,
                bucket="slide-visuals", destination_path=f"{session_id}/thumbnails/slide_{slide_number}.png"
            )

    return ProcessedSlide(
        slide_id=slide_id, slide_number=slide_number, svg_url=svg_url or "",
        original_width=slide_width_emu, original_height=slide_height_emu,
        thumbnail_url=thumbnail_url, shapes=extracted_shapes_data
    )


def get_slide_background_fill(slide) -> str:
    """
    Attempts to get the slide background fill color.
    Returns a hex color string (e.g., "#FFFFFF") or a default.
    Note: python-pptx has limitations in accessing complex background fills (gradients, pictures).
    This function will try to get solid fills.
    """
    try:
        fill = slide.background.fill
        if fill.type == 1:  # MSO_FILL.SOLID
            rgb = fill.fore_color.rgb
            return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
        # Handling for MSO_FILL.GRADIENT, MSO_FILL.PATTERN, MSO_FILL.PICTURE etc. is more complex
        # and often not fully exposed or easily convertible to a single SVG color.
    except Exception as e:
        logger.debug(f"Could not determine slide background color: {e}")
    return "#ffffff"  # Default to white


def create_minimal_svg(file_path: str, width_emu: int, height_emu: int):
    """Creates a minimal valid SVG file, e.g., if all generations fail."""
    EMU_PER_INCH = 914400
    DPI = 96
    width_px = int((width_emu / EMU_PER_INCH) * DPI)
    height_px = int((height_emu / EMU_PER_INCH) * DPI)

    svg_root = ET.Element('svg', xmlns='http://www.w3.org/2000/svg',
                          width=str(width_px), height=str(height_px),
                          viewBox=f'0 0 {width_px} {height_px}')
    ET.SubElement(svg_root, 'rect', width='100%',
                  height='100%', fill='#f0f0f0')
    text = ET.SubElement(svg_root, 'text', x='10', y='20', fill='red')
    text.text = "Error generating slide SVG"
    tree = ET.ElementTree(svg_root)
    ET.register_namespace('', 'http://www.w3.org/2000/svg')
    try:
        tree.write(file_path, encoding='utf-8', xml_declaration=True)
        logger.info(f"Created minimal fallback SVG: {file_path}")
    except Exception as e:
        logger.error(f"Failed to write minimal SVG {file_path}: {e}")

# Removed convert_slide_to_svg_using_libreoffice as it's replaced by _generate_svgs_for_all_slides_libreoffice


def extract_shapes(slide, slide_width_emu: int, slide_height_emu: int) -> List[SlideShape]:
    """
    Extract text and image shapes from a slide with their coordinates and styles.
    Coordinates are returned as percentages of slide dimensions.
    Dimensions are in EMU.
    """
    shapes_data = []
    # Ensure slide_width_emu and slide_height_emu are not zero to avoid DivisionByZero
    if slide_width_emu == 0 or slide_height_emu == 0:
        logger.warning(
            "Slide dimensions are zero, cannot calculate percentage-based shape coordinates.")
        return []

    for idx, shape in enumerate(slide.shapes):
        shape_left_emu = shape.left if shape.left is not None else 0
        shape_top_emu = shape.top if shape.top is not None else 0
        shape_width_emu = shape.width if shape.width is not None else 0
        shape_height_emu = shape.height if shape.height is not None else 0

        x_percent = (shape_left_emu / slide_width_emu) * 100
        y_percent = (shape_top_emu / slide_height_emu) * 100
        width_percent = (shape_width_emu / slide_width_emu) * 100
        height_percent = (shape_height_emu / slide_height_emu) * 100

        shape_obj_data = {
            "shape_id": str(uuid.uuid4()),
            "x_coordinate": x_percent,
            "y_coordinate": y_percent,
            "width": width_percent,
            "height": height_percent,
            "coordinates_unit": CoordinateUnit.PERCENTAGE,
            "reading_order": idx + 1,
            "original_text": None,  # Default
        }

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = shape.image
                image_bytes = image.blob
                image_base64_str = base64.b64encode(
                    image_bytes).decode('utf-8')
                shape_obj_data.update({
                    "shape_type": ShapeType.IMAGE,
                    "image_content_type": image.content_type,
                    "image_base64": image_base64_str
                })
                shapes_data.append(SlideShape(**shape_obj_data))
            except Exception as e:
                logger.warning(
                    f"Could not extract image data for shape {idx}: {e}")
                # Optionally, add a placeholder or skip
                continue

        elif shape.has_text_frame:
            text_frame = shape.text_frame
            full_text = text_frame.text.strip() if text_frame.text else ""
            if not full_text:
                continue

            shape_obj_data["original_text"] = full_text

            # Default style information
            font_size_pt = 12.0
            font_family = "Arial"
            font_weight = "normal"
            font_style = "normal"
            hex_color = "#000000"
            text_align_str = "LEFT"  # Default from PP_ALIGN
            vertical_anchor_str = "TOP"  # Default from MSO_VERTICAL_ANCHOR
            line_spacing_val = 1.0  # Default line spacing multiplier

            if text_frame.paragraphs:
                first_paragraph = text_frame.paragraphs[0]

                # Text Alignment
                if first_paragraph.alignment:
                    alignment_map = {
                        PP_ALIGN.LEFT: "LEFT", PP_ALIGN.CENTER: "CENTER",
                        PP_ALIGN.RIGHT: "RIGHT", PP_ALIGN.JUSTIFY: "JUSTIFY",
                        PP_ALIGN.DISTRIBUTE: "DISTRIBUTE", PP_ALIGN.THAI_DISTRIBUTE: "THAI_DISTRIBUTE"
                    }
                    text_align_str = alignment_map.get(
                        first_paragraph.alignment, "LEFT")

                # Line Spacing (complex, simplified here)
                # PPT stores line spacing in different ways (multiples of lines, points)
                # We simplify to a multiplier relative to font size.
                if first_paragraph.line_spacing is not None:
                    # Usually multiple of lines
                    if isinstance(first_paragraph.line_spacing, float):
                        line_spacing_val = first_paragraph.line_spacing
                    elif hasattr(first_paragraph.line_spacing, 'pt'):  # If in points
                        # Estimate based on a common default font size if run font size is not available
                        run_font_size_for_spacing = Pt(12)  # Default
                        if first_paragraph.runs and first_paragraph.runs[0].font and first_paragraph.runs[0].font.size:
                            run_font_size_for_spacing = first_paragraph.runs[0].font.size

                        if run_font_size_for_spacing and run_font_size_for_spacing.pt > 0:
                            line_spacing_val = first_paragraph.line_spacing.pt / run_font_size_for_spacing.pt
                        else:  # Fallback if font size is zero or unavailable
                            line_spacing_val = 1.15  # A common default multiplier

                # Font properties from the first run of the first paragraph
                if first_paragraph.runs:
                    first_run = first_paragraph.runs[0]
                    if first_run.font:
                        font = first_run.font
                        if font.size:
                            font_size_pt = font.size.pt
                        if font.name:
                            font_family = font.name
                        if font.bold:
                            font_weight = "bold"
                        if font.italic:
                            font_style = "italic"
                        if font.color.type == 1 and font.color.rgb:  # MSO_COLOR_TYPE.RGB
                            hex_color = f"#{font.color.rgb[0]:02x}{font.color.rgb[1]:02x}{font.color.rgb[2]:02x}"
                        # MSO_COLOR_TYPE.SCHEME is more complex, involves theme colors
                        # MSO_THEME_COLOR_INDEX
                        elif font.color.type == 2 and hasattr(font.color, 'theme_color'):
                            # This requires resolving theme colors, which is complex.
                            # For simplicity, we might ignore or use a default.
                            # logger.debug(f"Scheme color used: {font.color.theme_color}, brightness: {font.color.brightness}")
                            pass  # Placeholder for scheme color handling

            # Vertical Anchor for the text frame
            if text_frame.vertical_anchor:
                anchor_map = {
                    MSO_VERTICAL_ANCHOR.TOP: "TOP",
                    MSO_VERTICAL_ANCHOR.MIDDLE: "MIDDLE",
                    MSO_VERTICAL_ANCHOR.BOTTOM: "BOTTOM"
                }
                vertical_anchor_str = anchor_map.get(
                    text_frame.vertical_anchor, "TOP")

            shape_obj_data.update({
                "shape_type": ShapeType.TEXT,
                "font_size": font_size_pt,
                "font_family": font_family,
                "font_weight": font_weight,
                "font_style": font_style,
                "color": hex_color,
                "text_align": text_align_str,
                "vertical_anchor": vertical_anchor_str,
                "line_spacing": line_spacing_val,
            })
            shapes_data.append(SlideShape(**shape_obj_data))
    return shapes_data


def create_svg_from_slide(
    slide_shapes_data: List[SlideShape],
    file_path: str,
    width_emu: int,
    height_emu: int,
    slide_background_fill: str = "#ffffff"  # Added background fill parameter
) -> None:
    """
    Create an SVG representation of a PowerPoint slide using XML generation.
    Uses pre-extracted shapes_data.
    """
    EMU_PER_INCH = 914400
    DPI = 96
    POINTS_PER_INCH = 72

    width_px = max(1, int((width_emu / EMU_PER_INCH) * DPI))
    height_px = max(1, int((height_emu / EMU_PER_INCH) * DPI))

    svg_root = ET.Element('svg', xmlns='http://www.w3.org/2000/svg',
                          # Added for Pydantic v1 compatibility
                          xmlns_xlink='http://www.w3.org/1999/xlink',
                          width=str(width_px), height=str(height_px),
                          viewBox=f'0 0 {width_px} {height_px}')

    # Slide Background
    background = ET.SubElement(
        svg_root, 'rect', width='100%', height='100%', fill=slide_background_fill)

    # Optional: Slide Border (can be made configurable)
    # border = ET.SubElement(svg_root, 'rect', x='0', y='0', width=str(width_px), height=str(height_px),
    #                       fill='none', stroke='#e0e0e0', stroke_width='1')

    for shape_data in slide_shapes_data:
        # Calculate pixel values from percentage data
        x_px_shape = int((shape_data.x_coordinate / 100) * width_px)
        y_px_shape = int((shape_data.y_coordinate / 100) * height_px)
        # Ensure width is at least 1px
        w_px_shape = max(1, int((shape_data.width / 100) * width_px))
        # Ensure height is at least 1px
        h_px_shape = max(1, int((shape_data.height / 100) * height_px))

        if shape_data.shape_type == ShapeType.IMAGE and shape_data.image_base64:
            image_element = ET.SubElement(svg_root, 'image',
                                          x=str(x_px_shape), y=str(y_px_shape),
                                          width=str(w_px_shape), height=str(h_px_shape),
                                          preserveAspectRatio="xMidYMid meet")  # Added preserveAspectRatio
            image_element.set('{http://www.w3.org/1999/xlink}href',
                              f"data:{shape_data.image_content_type};base64,{shape_data.image_base64}")

        elif shape_data.shape_type == ShapeType.TEXT and shape_data.original_text:
            text_g = ET.SubElement(svg_root, 'g', id=shape_data.shape_id,
                                   transform=f"translate({x_px_shape},{y_px_shape})")

            # Optional: Add a bounding box for the text container for visualization/debugging
            # text_bg_rect = ET.SubElement(text_g, 'rect',
            #                              x="0", y="0",
            #                              width=str(w_px_shape), height=str(h_px_shape),
            #                              fill="rgba(200,200,200,0.1)", stroke="#cccccc", stroke_width="0.5")

            text_element = ET.SubElement(text_g, 'text',
                                         font_family=shape_data.font_family or "Arial",
                                         fill=shape_data.color or "#000000")

            font_size_px = max(
                1, int((shape_data.font_size or 12.0) * (DPI / POINTS_PER_INCH)))
            text_element.set('font-size', str(font_size_px))

            if shape_data.font_weight == "bold":
                text_element.set('font-weight', 'bold')
            if shape_data.font_style == "italic":
                text_element.set('font-style', 'italic')

            # Text Alignment (text-anchor)
            text_anchor = "start"  # Default for LTR languages
            # Adjust x for text-anchor. Margin for padding within the shape.
            text_x_offset = 5
            if shape_data.text_align == "CENTER":
                text_anchor = "middle"
                text_x_offset = w_px_shape / 2
            elif shape_data.text_align == "RIGHT":
                text_anchor = "end"
                text_x_offset = w_px_shape - 5
            text_element.set('text-anchor', text_anchor)

            # Vertical Alignment (dominant-baseline and y position of first tspan)
            # This is tricky in SVG. We try to approximate.
            # 'dy' on tspans handles line spacing. The initial 'y' sets the first line's position.

            # Estimate first line's Y based on vertical anchor
            # This requires careful adjustment. dominant-baseline is key.
            # Default: top-aligned-ish (cap-height)
            first_line_y_px = font_size_px

            if shape_data.vertical_anchor == "TOP":
                text_element.set('dominant-baseline',
                                 'text-before-edge')  # or 'hanging'
                first_line_y_px = 5  # Small padding from top
            elif shape_data.vertical_anchor == "MIDDLE":
                text_element.set('dominant-baseline', 'central')  # or 'middle'
                first_line_y_px = h_px_shape / 2
            elif shape_data.vertical_anchor == "BOTTOM":
                text_element.set('dominant-baseline', 'text-after-edge')
                first_line_y_px = h_px_shape - 5  # Small padding from bottom
            else:  # Default to hanging/text-before-edge for better top alignment
                text_element.set('dominant-baseline', 'hanging')
                first_line_y_px = 5

            lines = shape_data.original_text.splitlines()
            if not lines:
                # Ensure at least one tspan if text is empty but shape exists
                lines = [' ']

            line_spacing_multiplier = shape_data.line_spacing if (
                shape_data.line_spacing and shape_data.line_spacing > 0) else 1.15
            # Actual line height in pixels for use with 'dy'
            actual_line_height_px = font_size_px * line_spacing_multiplier

            for i, line_text in enumerate(lines):
                tspan = ET.SubElement(text_element, 'tspan',
                                      x=str(text_x_offset))
                # Use non-breaking space for empty lines
                tspan.text = line_text if line_text.strip() else ' '

                if i == 0:
                    tspan.set('y', str(first_line_y_px))
                else:
                    # Use 'px' for clarity, though unitless often works for dy
                    tspan.set('dy', f"{actual_line_height_px}px")

    tree = ET.ElementTree(svg_root)
    ET.register_namespace('', 'http://www.w3.org/2000/svg')
    ET.register_namespace('xlink', 'http://www.w3.org/1999/xlink')
    try:
        tree.write(file_path, encoding='utf-8', xml_declaration=True)
    except Exception as e:
        logger.error(f"Error writing SVG file {file_path}: {e}")


def create_thumbnail_from_slide_pil(
    slide,  # python-pptx slide object
    shapes_data: List[SlideShape],  # Pre-extracted shapes data
    file_path: str,
    slide_width_emu: int,
    slide_height_emu: int,
    thumbnail_width_px: int = 250
) -> None:
    """
    Create a thumbnail image for the slide using PIL, using extracted shapes data.
    """
    if slide_width_emu == 0 or slide_height_emu == 0:
        logger.warning("Cannot generate thumbnail, slide dimensions are zero.")
        # Create a tiny placeholder image
        img = Image.new('RGB', (thumbnail_width_px,
                        thumbnail_width_px // 2), color=(230, 230, 230))
        draw = ImageDraw.Draw(img)
        draw.text((10, 10), "Error: Zero slide dimensions", fill=(255, 0, 0))
        img.save(file_path)
        return

    aspect_ratio = slide_height_emu / slide_width_emu
    thumbnail_height_px = int(aspect_ratio * thumbnail_width_px)

    # Use slide background color if available from slide object (simplified)
    # For more accuracy, this could be passed or enhanced.
    slide_bg_hex = get_slide_background_fill(slide)
    try:
        from PIL import ImageColor
        image_bg_color = ImageColor.getrgb(slide_bg_hex)
    except ValueError:
        image_bg_color = (255, 255, 255)  # Default white

    image = Image.new('RGB', (thumbnail_width_px,
                      thumbnail_height_px), color=image_bg_color)
    draw = ImageDraw.Draw(image)

    # Draw slide border
    draw.rectangle(
        [(0, 0), (thumbnail_width_px-1, thumbnail_height_px-1)], outline=(200, 200, 200))

    default_font = None
    try:
        default_font = ImageFont.truetype("arial.ttf", 10)
    except IOError:
        try:
            default_font = ImageFont.load_default()  # Fallback to PIL default bitmap font
        except IOError:
            logger.warning(
                "Thumbnail font not found, text rendering in thumbnails will be basic lines.")

    for shape_data in shapes_data:
        # Calculate thumbnail coordinates from percentage
        x_thumb = int((shape_data.x_coordinate / 100) * thumbnail_width_px)
        y_thumb = int((shape_data.y_coordinate / 100) * thumbnail_height_px)
        w_thumb = int((shape_data.width / 100) * thumbnail_width_px)
        h_thumb = int((shape_data.height / 100) * thumbnail_height_px)

        if shape_data.shape_type == ShapeType.IMAGE and shape_data.image_base64:
            try:
                img_data = base64.b64decode(shape_data.image_base64)
                from io import BytesIO
                img_io = BytesIO(img_data)
                shape_img = Image.open(img_io)
                shape_img = shape_img.resize(
                    (w_thumb, h_thumb), Image.Resampling.LANCZOS)
                image.paste(shape_img, (x_thumb, y_thumb),
                            mask=shape_img.convert("RGBA"))
            except Exception as e:
                logger.warning(
                    f"Could not render image in thumbnail: {e}. Drawing placeholder.")
                draw.rectangle([(x_thumb, y_thumb), (x_thumb + w_thumb, y_thumb + h_thumb)],
                               fill=(200, 220, 255), outline=(150, 180, 230))
                if default_font:
                    draw.text((x_thumb + 2, y_thumb + 2),
                              "[img]", font=default_font, fill=(0, 0, 0))

        elif shape_data.shape_type == ShapeType.TEXT and shape_data.original_text:
            # Draw text block placeholder
            draw.rectangle([(x_thumb, y_thumb), (x_thumb + w_thumb, y_thumb + h_thumb)],
                           fill=(240, 240, 240), outline=(180, 180, 180))

            if default_font:
                text_preview = (shape_data.original_text[:20] + '...') if len(
                    shape_data.original_text) > 20 else shape_data.original_text
                # Attempt to draw text, handle very small boxes
                if w_thumb > 5 and h_thumb > 5:
                    draw.text((x_thumb + 2, y_thumb + 2), text_preview,
                              font=default_font, fill=(50, 50, 50))
            else:  # Fallback if font is not available
                draw.line([(x_thumb+2, y_thumb+h_thumb//2), (x_thumb +
                          w_thumb-2, y_thumb+h_thumb//2)], fill=(100, 100, 100))
    try:
        image.save(file_path)
    except Exception as e:
        logger.error(f"Error saving thumbnail {file_path}: {e}")
        # Attempt to save a minimal error image
        try:
            error_img = Image.new('RGB', (50, 20), color=(255, 0, 0))
            ImageDraw.Draw(error_img).text((2, 2), "ERR", fill=(255, 255, 255))
            error_img.save(file_path)
        except:
            pass  # Give up if even error image fails
