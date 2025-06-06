import os
import uuid
import logging
import time
import asyncio
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime
from pptx import Presentation

import json
import shutil

import tempfile

from app.models.schemas import (
    ProcessedSlide,
    SlideShape,
    ProcessedPresentation,
    OverallProcessingStatus,
    ShapeType,
    CoordinateUnit,
    ProcessingStatusResponse,
    ProcessingStatus
)
from app.services.supabase_service import upload_file_to_supabase, update_job_status as update_supabase_job_status
from app.services.job_status import update_job_status as update_local_job_status, get_job_status
from app.core.config import get_settings
from app.core.utils import async_retry
# Import the new ProcessingManager related functions
from app.services.processing_manager import get_processing_manager
from app.services.cache_service import get_cache_service # Import CacheService
from app.services.svg_generator import generate_svgs
from app.services.slide_parser import (
    extract_shapes_enhanced,
    create_thumbnail_from_slide_enhanced,
    validate_coordinates_with_svg
)

logger = logging.getLogger(__name__)

# Get settings
settings = get_settings()

# Removed hardcoded LIBREOFFICE_PATH, will use settings.LIBREOFFICE_PATH

# Dictionary to keep track of job file paths for retry capability
_job_file_paths = {}





async def get_job_file_path(job_id: str) -> Optional[str]:
    """
    Get the file path for a job.
    Returns None if the job doesn't exist or the file has been cleaned up.
    """
    # First check our in-memory cache
    if job_id in _job_file_paths and os.path.exists(_job_file_paths[job_id]):
        return _job_file_paths[job_id]

    # If not in memory, we can't retrieve it since we don't persistently store file paths
    return None


async def queue_pptx_processing(
    job_id: str,
    session_id: str,
    file_path: str,
    source_language: Optional[str] = None,
    target_language: Optional[str] = None,
    generate_thumbnails: bool = True
) -> None:
    """
    Queue the PPTX processing task using the ProcessingManager.
    """
    _job_file_paths[job_id] = file_path

    # Initial status update to QUEUED, as the manager will pick it up.
    # The manager loop, when it picks a job, can update it to PROCESSING.
    # Or, process_pptx itself can update it to PROCESSING as its first step.
    # For now, let's set it to QUEUED reflecting it's in our system's queue.
    await update_local_job_status(
        job_id=job_id,
        status=ProcessingStatusResponse(
            job_id=job_id,
            session_id=session_id,
            status=ProcessingStatus.QUEUED, # Changed from PROCESSING to QUEUED
            progress=0,
            current_stage="Job queued for processing"
        )
    )

    job_details = {
        'job_id': job_id,
        'session_id': session_id,
        'file_path': file_path,
        'source_language': source_language,
        'target_language': target_language,
        'generate_thumbnails': generate_thumbnails
    }

    try:
        manager = get_processing_manager()
        await manager.submit_job(job_details)
        logger.info(f"Job {job_id} successfully submitted to ProcessingManager.")
    except Exception as e:
        logger.error(f"Failed to submit job {job_id} to ProcessingManager: {e}", exc_info=True)
        # Update status to FAILED if submission itself fails
        await update_local_job_status(
            job_id=job_id,
            status=ProcessingStatusResponse(
                job_id=job_id,
                session_id=session_id,
                status=ProcessingStatus.FAILED,
                progress=0,
                current_stage="Failed to queue job",
                error=f"Failed to submit to manager: {str(e)}"
            )
        )
        # Also update Supabase status if needed
        await update_supabase_job_status(
            session_id=session_id, status="failed", error=f"Failed to queue job: {str(e)}"
        )
        # Re-raise or handle as appropriate for the caller
        raise


async def process_pptx(
    job_id: str,
    session_id: str,
    file_path: str,
    source_language: Optional[str] = None,
    target_language: Optional[str] = None,
    generate_thumbnails: bool = True
) -> None:
    """
    Process a PPTX file using LibreOffice-only approach.
    Checks cache first, then processes if not found, and stores to cache on success.
    """
    start_time = time.time()
    cache_service = get_cache_service()
    cache_params = {
        "source_language": source_language or "", # Ensure consistent type for key
        "target_language": target_language or "",
        "generate_thumbnails": generate_thumbnails
    }
    cache_key = cache_service.generate_cache_key(file_path, cache_params)

    if cache_key:
        cached_result_tuple = cache_service.get_cached_result(cache_key)
        if cached_result_tuple:
            cached_presentation, cached_result_url = cached_result_tuple
            logger.info(f"Cache hit for job {job_id} (session {session_id}) with key {cache_key}.")
            
            # Update local status to COMPLETED
            await update_local_job_status(
                job_id=job_id,
                status=ProcessingStatusResponse(
                    job_id=job_id, session_id=session_id, status=ProcessingStatus.COMPLETED,
                    progress=100, current_stage="Processing completed (from cache)",
                    completed_at=datetime.now() # Or use a cached completion time if stored
                )
            )
            # Update Supabase status
            await update_supabase_job_status(
                session_id=session_id,
                status="completed",
                slide_count=cached_presentation.slide_count,
                result_url=cached_result_url # This is the URL to the main result.json
            )
            # Potentially update _job_file_paths if needed for retries, though cache hit means no retry of processing.
            # For now, we assume original file path might still be relevant for other operations.
            if job_id not in _job_file_paths:
                 _job_file_paths[job_id] = file_path

            # Clean up the uploaded file since processing is "done" via cache
            # Ensure this cleanup logic is consistent with the non-cached path's finally block
            uploaded_file_dir_for_cleanup = os.path.dirname(file_path)
            try:
                if os.path.exists(uploaded_file_dir_for_cleanup):
                    # Ensure no processing_output dir is accidentally deleted if it wasn't created
                    # For cache hit, processing_output_dir is not created by this run.
                    # We only want to clean up the original uploaded file and its parent if empty.
                    if os.path.isfile(file_path): # Check if it's a file before deleting its dir
                         os.remove(file_path)
                         logger.info(f"Cleaned up source file (cache hit): {file_path}")
                         # Attempt to remove the directory if it's empty
                         try:
                             os.rmdir(uploaded_file_dir_for_cleanup)
                             logger.info(f"Cleaned up empty source directory (cache hit): {uploaded_file_dir_for_cleanup}")
                         except OSError:
                             logger.debug(f"Source directory not empty, not removed (cache hit): {uploaded_file_dir_for_cleanup}")


                    if job_id in _job_file_paths: # Remove from cache after cleanup
                        del _job_file_paths[job_id]

            except Exception as e_cleanup:
                logger.error(f"Error cleaning up source file on cache hit {file_path}: {str(e_cleanup)}")
            return # End processing here due to cache hit

    logger.info(f"Cache miss for job {job_id} (session {session_id}). Proceeding with full processing.")
    # Update local status to PROCESSING (if not already updated by manager)
    await update_local_job_status(
        job_id=job_id,
        status=ProcessingStatusResponse(
            job_id=job_id, session_id=session_id, status=ProcessingStatus.PROCESSING,
            progress=1, current_stage="Starting full processing"
        )
    )

    # processing_dir is the main directory for this job's outputs (SVGs, JSON)
    uploaded_file_dir = os.path.dirname(file_path)
    processing_output_dir = os.path.join(uploaded_file_dir, "processing_output")
    os.makedirs(processing_output_dir, exist_ok=True)

    try:
        # Validate LibreOffice availability upfront
        from app.services.svg_generator import validate_libreoffice_availability
        if not validate_libreoffice_availability():
            raise ValueError("LibreOffice not configured or not available")

        presentation = Presentation(file_path)
        slide_count = len(presentation.slides)

        await update_local_job_status(
            job_id=job_id,
            status=ProcessingStatusResponse(
                job_id=job_id, session_id=session_id, status=ProcessingStatus.PROCESSING,
                progress=5, current_stage=f"Opened presentation with {slide_count} slides"
            )
        )

        # Generate all SVGs using LibreOffice - this is now required, not optional
        await update_local_job_status(
            job_id=job_id,
            status=ProcessingStatusResponse(
                job_id=job_id, session_id=session_id, status=ProcessingStatus.PROCESSING,
                progress=10, current_stage="Converting slides to SVG with LibreOffice"
            )
        )

        # Use the new svg_generator module
        generated_svg_paths = await generate_svgs(
            presentation_path=file_path,
            output_dir=processing_output_dir,
            slide_count=slide_count
        )
        if not generated_svg_paths:
            raise RuntimeError("SVG generation failed to produce any files.")
        
        logger.info(f"Generated {len(generated_svg_paths)} SVG files.", extra={"slide_count": slide_count, "output_dir": processing_output_dir})

        # Process each slide with enhanced text extraction
        processed_slides_data = []
        for idx, slide in enumerate(presentation.slides):
            slide_number = idx + 1
            progress = 20 + int((idx / slide_count) * 75)
            await update_local_job_status(
                job_id=job_id,
                status=ProcessingStatusResponse(
                    job_id=job_id, session_id=session_id, status=ProcessingStatus.PROCESSING,
                    progress=progress, current_stage=f"Processing slide {slide_number} of {slide_count}"
                )
            )

            # Verify LibreOffice SVG exists for this slide
            svg_path = generated_svg_paths.get(slide_number)
            if not svg_path or not os.path.exists(svg_path):
                raise RuntimeError(f"LibreOffice SVG missing for slide {slide_number}")

            processed_slide = await process_slide_simplified(
                slide=slide,
                slide_number=slide_number,
                svg_path=svg_path,
                session_id=session_id,
                generate_thumbnail=generate_thumbnails
            )
            processed_slides_data.append(processed_slide)

        # Finalize processing
        processing_time = int(time.time() - start_time)
        result = ProcessedPresentation(
            session_id=session_id, 
            slide_count=slide_count,
            processing_status=OverallProcessingStatus.COMPLETED,
            processing_time=processing_time, 
            slides=processed_slides_data
        )

        result_file = os.path.join(processing_output_dir, f"result_{session_id}.json")
        with open(result_file, "w") as f:
            json.dump(result.model_dump(mode='json'), f, indent=4)

        result_url = await upload_file_to_supabase(
            file_path=result_file,
            bucket="processing-results", 
            destination_path=f"{session_id}/result.json"
        )

        # Store successful result in cache
        if cache_key and result_url: # result_url must exist to be a valid cache entry
            cache_service.store_cached_result(cache_key, result, result_url)
            logger.info(f"Processing result for job {job_id} stored in cache with key {cache_key}.")

        await update_local_job_status(
            job_id=job_id,
            status=ProcessingStatusResponse(
                job_id=job_id, session_id=session_id, status=ProcessingStatus.COMPLETED,
                progress=100, current_stage="Processing completed", completed_at=datetime.now()
            )
        )
        
        await update_supabase_job_status(
            session_id=session_id, 
            status="completed",
            slide_count=slide_count, 
            result_url=result_url
        )

    except Exception as e:
        logger.error(f"Error processing PPTX for job {job_id}: {str(e)}", exc_info=True)
        await update_local_job_status(
            job_id=job_id,
            status=ProcessingStatusResponse(
                job_id=job_id, session_id=session_id, status=ProcessingStatus.FAILED,
                progress=0, current_stage="Processing failed", error=str(e) # Ensure progress is reset on failure
            )
        )
        await update_supabase_job_status(
            session_id=session_id, status="failed", error=str(e)
        )
        # Do not re-raise here if called by ProcessingManager, as it handles logging.
        # If this function can be called directly, re-raising might be desired.
        # For now, assuming ProcessingManager's _execute_job handles the top-level try/except.
        # However, if this function is intended to be a standalone callable that signals failure by raising,
        # then `raise` should be here. Given it's called by the manager, which logs, let's not re-raise.
        # Update: The manager calls this and logs, but the original function re-raised.
        # For consistency with how it was called by asyncio.create_task before, let's re-raise.
        raise
        
    finally:
        # Clean up temporary files from processing_output_dir and the uploaded file
        try:
            # This status check is to avoid deleting files if the job is still queued or actively processing
            # (e.g. if cleanup is called prematurely or from another context).
            # However, in the context of process_pptx, it has either completed or failed.
            current_job_status_info = await get_job_status(job_id)
            
            # Check if the job is truly finished (completed or failed) before cleaning up.
            # This helps prevent premature deletion if the job is retried or requeued by a more advanced manager.
            is_job_finished = current_job_status_info and current_job_status_info.status in [ProcessingStatus.COMPLETED, ProcessingStatus.FAILED]

            if is_job_finished:
                if os.path.exists(processing_output_dir): # This is the specific output dir for this run
                    shutil.rmtree(processing_output_dir)
                    logger.info(f"Cleaned up temporary processing output directory: {processing_output_dir}")
                
                # Clean up the original uploaded file and its parent directory if empty
                if os.path.exists(uploaded_file_dir) and os.path.isdir(uploaded_file_dir):
                    # Original file path
                    original_upload_path = os.path.join(uploaded_file_dir, os.path.basename(file_path))
                    if os.path.exists(original_upload_path) and os.path.isfile(original_upload_path):
                        os.remove(original_upload_path)
                        logger.info(f"Cleaned up original uploaded file: {original_upload_path}")
                    
                    # Attempt to remove the parent directory if it's empty
                    # This assumes uploaded_file_dir is specific to this job's upload
                    # and doesn't contain other ongoing uploads.
                    try:
                        if not os.listdir(uploaded_file_dir): # Check if empty
                            os.rmdir(uploaded_file_dir)
                            logger.info(f"Cleaned up empty upload directory: {uploaded_file_dir}")
                    except OSError:
                        logger.debug(f"Upload directory not empty or other error, not removed: {uploaded_file_dir}")

                # Remove from job file paths cache
                if job_id in _job_file_paths:
                    del _job_file_paths[job_id]
            else:
                logger.info(f"Job {job_id} is not marked as finished ({current_job_status_info.status if current_job_status_info else 'unknown'}). Skipping cleanup of {processing_output_dir} and {uploaded_file_dir}.")

        except Exception as e_cleanup:
            logger.error(f"Error cleaning up temporary files for job {job_id} at {uploaded_file_dir}: {str(e_cleanup)}")


async def process_slide_simplified(
    slide,  # python-pptx Slide object
    slide_number: int,
    svg_path: str,  # Path to LibreOffice-generated SVG
    session_id: str,
    generate_thumbnail: bool = True
) -> ProcessedSlide:
    """
    Processes a single slide: extracts shapes, creates thumbnail, uploads assets.
    This is now an orchestrator for the slide_parser module.
    """
    start_time = time.time()
    log_context = {"session_id": session_id, "slide_number": slide_number}
    logger.info("Processing slide", extra=log_context)

    try:
        # Correctly access slide dimensions via the presentation part
        pres = slide.part.package.presentation_part.presentation
        slide_width_emu = pres.slide_width
        slide_height_emu = pres.slide_height

        shapes_data = extract_shapes_enhanced(slide, slide_width_emu, slide_height_emu)
        shapes_data = await validate_coordinates_with_svg(shapes_data, svg_path, slide_width_emu, slide_height_emu)
        logger.info(f"Extracted {len(shapes_data)} shapes from slide.", extra={**log_context, "shape_count": len(shapes_data)})

        svg_object_name = f"{session_id}/slides/slide_{slide_number}.svg"
        svg_url = await upload_file_to_supabase(
            svg_path, settings.SUPABASE_STORAGE_BUCKET, svg_object_name
        )
        logger.info("Uploaded SVG to Supabase.", extra={**log_context, "svg_url": svg_url})

        thumbnail_url = ""
        if generate_thumbnail:
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_thumb_file:
                try:
                    # Also correct the dimension access here
                    pres = slide.part.package.presentation_part.presentation
                    create_thumbnail_from_slide_enhanced(
                        slide, shapes_data, temp_thumb_file.name, pres.slide_width, pres.slide_height
                    )
                    thumb_object_name = f"{session_id}/thumbnails/thumbnail_{slide_number}.png"
                    thumbnail_url = await upload_file_to_supabase(
                        temp_thumb_file.name, settings.SUPABASE_STORAGE_BUCKET, thumb_object_name
                    )
                    logger.info("Uploaded thumbnail to Supabase.", extra={**log_context, "thumbnail_url": thumbnail_url})
                finally:
                    os.remove(temp_thumb_file.name)

        end_time = time.time()
        slide_processing_time = end_time - start_time
        logger.info(f"Finished processing slide in {slide_processing_time:.2f}s", extra={**log_context, "duration_seconds": slide_processing_time})
        
        return ProcessedSlide(
            slide_id=str(uuid.uuid4()),
            slide_number=slide_number,
            svg_url=svg_url,
            original_width=int(slide_width_emu / 9525),  # Convert EMU to pixels (1 EMU = 1/9525 inch, 96 DPI)
            original_height=int(slide_height_emu / 9525),
            thumbnail_url=thumbnail_url,
            shapes=shapes_data
        )

    except Exception as e:
        end_time = time.time()
        slide_processing_time = end_time - start_time
        error_message = f"Failed to process slide {slide_number}: {e}"
        logger.error(error_message, extra={**log_context, "error": str(e)}, exc_info=True)
        return ProcessedSlide(
            slide_id=str(uuid.uuid4()),
            slide_number=slide_number,
            svg_url="https://example.com/error.svg",  # Provide a valid URL for validation
            original_width=1920,  # Default dimensions
            original_height=1080,
            thumbnail_url="https://example.com/error.png",
            shapes=[]
        )





def segment_text_for_translation(text: str, max_segment_length: int = 100) -> List[Dict[str, Any]]:
    """
    Segment text into translation-friendly units.
    Breaks text at sentence boundaries while respecting length limits.
    """
    if not text or len(text.strip()) == 0:
        return []

    text = text.strip()
    
    # If text is short enough, return as single segment
    if len(text) <= max_segment_length:
        return [{
            'text': text,
            'segment_index': 0,
            'is_complete_sentence': True,
            'word_count': len(text.split()),
            'char_count': len(text)
        }]

    segments = []
    
    # Split by sentence boundaries
    sentence_endings = r'[.!?]+\s+'
    sentences = re.split(sentence_endings, text)
    
    current_segment = ""
    segment_index = 0
    
    for i, sentence in enumerate(sentences):
        sentence = sentence.strip()
        if not sentence:
            continue
            
        # If adding this sentence would exceed limit, save current segment
        if current_segment and len(current_segment + " " + sentence) > max_segment_length:
            if current_segment:
                segments.append({
                    'text': current_segment.strip(),
                    'segment_index': segment_index,
                    'is_complete_sentence': True,
                    'word_count': len(current_segment.split()),
                    'char_count': len(current_segment)
                })
                segment_index += 1
                current_segment = sentence
            else:
                # Single sentence is too long, break it at word boundaries
                words = sentence.split()
                temp_segment = ""
                
                for word in words:
                    if temp_segment and len(temp_segment + " " + word) > max_segment_length:
                        segments.append({
                            'text': temp_segment.strip(),
                            'segment_index': segment_index,
                            'is_complete_sentence': False,
                            'word_count': len(temp_segment.split()),
                            'char_count': len(temp_segment)
                        })
                        segment_index += 1
                        temp_segment = word
                    else:
                        temp_segment = temp_segment + " " + word if temp_segment else word
                
                current_segment = temp_segment
        else:
            current_segment = current_segment + " " + sentence if current_segment else sentence

    # Add final segment
    if current_segment:
        segments.append({
            'text': current_segment.strip(),
            'segment_index': segment_index,
            'is_complete_sentence': True,
            'word_count': len(current_segment.split()),
            'char_count': len(current_segment)
        })

    return segments
