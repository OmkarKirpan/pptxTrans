import os
import uuid
import logging
import time
import asyncio
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import json
import shutil
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Emu, Pt
import base64
from pptx.enum.shapes import MSO_SHAPE_TYPE
import xml.etree.ElementTree as ET
import re

from app.models.schemas import SlideShape, ShapeType, CoordinateUnit

try:
    from fuzzywuzzy import fuzz
except ImportError:
    fuzz = None

logger = logging.getLogger(__name__)

def get_slide_background_fill(slide) -> str:
    """
    Attempts to get the slide background fill color.
    Returns a hex color string (e.g., "#FFFFFF") or a default.
    Note: python-pptx has limitations in accessing complex background fills (gradients, pictures).
    This function will try to get solid fills.
    """
    try:
        fill = slide.background.fill
        if fill.type == 1:  # MSO_FILL.SOLID
            rgb = fill.fore_color.rgb
            return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
        # Handling for MSO_FILL.GRADIENT, MSO_FILL.PATTERN, MSO_FILL.PICTURE etc. is more complex
        # and often not fully exposed or easily convertible to a single SVG color.
    except Exception as e:
        logger.debug(f"Could not determine slide background color: {e}")
    return "#ffffff"  # Default to white

def extract_shapes_enhanced(slide, slide_width_emu: int, slide_height_emu: int) -> List[SlideShape]:
    shapes_data: List[SlideShape] = []
    for shape in slide.shapes:
        if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            table_left = shape.left if shape.left else 0
            table_top = shape.top if shape.top else 0
            col_widths = [col.width for col in table.columns]
            row_heights = [row.height for row in table.rows]
            
            running_top = table_top
            for r, row in enumerate(table.rows):
                running_left = table_left
                for c, cell in enumerate(row.cells):
                    if cell.text and cell.text.strip():
                        text = cell.text
                        shapes_data.append(SlideShape(
                            shape_id=str(uuid.uuid4()), 
                            shape_type=ShapeType.TABLE_CELL, 
                            original_text=text,
                            x_coordinate=running_left, 
                            y_coordinate=running_top, 
                            width=col_widths[c], 
                            height=row_heights[r],
                            coordinates_unit=CoordinateUnit.EMU,
                            is_title=False, 
                            is_subtitle=False, 
                            placeholder_type="TABLE_CELL",
                            text_length=len(text), 
                            word_count=len(text.split()),
                            translation_priority=5,
                            validation_status="unvalidated",
                            validation_details="Coordinates have not been validated against SVG output."
                        ))
                    running_left += col_widths[c]
                running_top += row_heights[r]
            continue

        if not shape.has_text_frame or not hasattr(shape, 'text') or not shape.text or not shape.text.strip():
            continue

        text_frame = shape.text_frame
        text = text_frame.text
        is_title = shape.is_placeholder and (shape.placeholder_format.type in [1, 101])
        is_subtitle = shape.is_placeholder and shape.placeholder_format.type == 2
        
        shapes_data.append(SlideShape(
            shape_id=str(uuid.uuid4()), 
            shape_type=ShapeType.TEXT, 
            original_text=text,
            x_coordinate=shape.left if shape.left else 0, 
            y_coordinate=shape.top if shape.top else 0,
            width=shape.width if shape.width else 0, 
            height=shape.height if shape.height else 0,
            coordinates_unit=CoordinateUnit.EMU,
            is_title=is_title, 
            is_subtitle=is_subtitle,
            placeholder_type=shape.placeholder_format.type.name if shape.is_placeholder else "NONE",
            text_length=len(text), 
            word_count=len(text.split()),
            translation_priority=10 if is_title else (8 if is_subtitle else 5),
            validation_status="unvalidated"
        ))
    return shapes_data

def create_thumbnail_from_slide_enhanced(slide, shapes_data: List[SlideShape], file_path: str, slide_width_emu: int, slide_height_emu: int, thumbnail_width_px: int = 400):
    if slide_width_emu == 0 or slide_height_emu == 0: return
    aspect_ratio = slide_height_emu / slide_width_emu
    thumbnail_height_px = int(thumbnail_width_px * aspect_ratio)
    img = Image.new('RGB', (thumbnail_width_px, thumbnail_height_px), color=get_slide_background_fill(slide))
    draw = ImageDraw.Draw(img)
    scale_x, scale_y = thumbnail_width_px / slide_width_emu, thumbnail_height_px / slide_height_emu
    for shape in shapes_data:
        x1, y1 = int(shape.x_coordinate * scale_x), int(shape.y_coordinate * scale_y)
        x2, y2 = int((shape.x_coordinate + shape.width) * scale_x), int((shape.y_coordinate + shape.height) * scale_y)
        draw.rectangle([x1, y1, x2, y2], outline="lightgrey", fill="#F0F0F0")
    img.save(file_path, 'PNG')

def _mark_shapes_as_unvalidated(shapes_data: List[SlideShape], reason: str) -> List[SlideShape]:
    logger.warning(f"Coordinate validation skipped: {reason}", extra={"reason": reason})
    for shape in shapes_data:
        shape.validation_status = "unvalidated"
        shape.validation_details = reason
    return shapes_data

def _extract_svg_dimensions(root) -> Dict[str, Any]:
    """Extract SVG viewport dimensions and viewBox information."""
    try:
        # Get SVG element attributes
        width = root.get('width', '0')
        height = root.get('height', '0')
        viewbox = root.get('viewBox', '0 0 0 0')
        
        # Parse viewBox if available
        viewbox_parts = viewbox.strip().split()
        if len(viewbox_parts) == 4:
            viewbox_x, viewbox_y, viewbox_width, viewbox_height = map(float, viewbox_parts)
        else:
            viewbox_x, viewbox_y, viewbox_width, viewbox_height = 0, 0, 0, 0
        
        # Extract numeric values from width/height (remove units like 'px', 'pt')
        import re
        width_match = re.search(r'(\d+(?:\.\d+)?)', str(width))
        height_match = re.search(r'(\d+(?:\.\d+)?)', str(height))
        
        svg_width = float(width_match.group(1)) if width_match else viewbox_width
        svg_height = float(height_match.group(1)) if height_match else viewbox_height
        
        return {
            'svg_width': svg_width,
            'svg_height': svg_height,
            'viewbox_x': viewbox_x,
            'viewbox_y': viewbox_y,
            'viewbox_width': viewbox_width,
            'viewbox_height': viewbox_height
        }
    except Exception as e:
        logger.warning(f"Could not extract SVG dimensions: {e}")
        return {}

def _calculate_coordinate_transform(svg_info, slide_width_px, slide_height_px) -> Dict[str, float]:
    """Calculate transformation parameters between slide coordinates and SVG coordinates."""
    if not svg_info or 'svg_width' not in svg_info:
        return {}
    
    try:
        # Use viewBox dimensions if available, otherwise use SVG dimensions
        svg_width = svg_info.get('viewbox_width') or svg_info.get('svg_width', 0)
        svg_height = svg_info.get('viewbox_height') or svg_info.get('svg_height', 0)
        
        if svg_width == 0 or svg_height == 0:
            return {}
        
        # Calculate scale factors
        scale_x = svg_width / slide_width_px if slide_width_px > 0 else 1.0
        scale_y = svg_height / slide_height_px if slide_height_px > 0 else 1.0
        
        # Account for viewBox offset
        offset_x = svg_info.get('viewbox_x', 0)
        offset_y = svg_info.get('viewbox_y', 0)
        
        return {
            'scale_x': scale_x,
            'scale_y': scale_y,
            'offset_x': offset_x,
            'offset_y': offset_y,
            'svg_width': svg_width,
            'svg_height': svg_height
        }
    except Exception as e:
        logger.warning(f"Could not calculate coordinate transform: {e}")
        return {}

def _extract_svg_text_elements(root, namespaces) -> List[Dict[str, Any]]:
    """Extract text elements and their positions from SVG."""
    text_elements = []
    
    try:
        # Find all text elements (both <text> and <tspan>)
        for text_elem in root.iter():
            if text_elem.tag.endswith('}text') or text_elem.tag == 'text':
                # Extract text content
                text_content = ''.join(text_elem.itertext()).strip()
                if not text_content:
                    continue
                
                # Extract position attributes
                x = float(text_elem.get('x', 0))
                y = float(text_elem.get('y', 0))
                
                # Extract styling information
                font_size = text_elem.get('font-size', '')
                font_family = text_elem.get('font-family', '')
                
                text_elements.append({
                    'text': text_content,
                    'x': x,
                    'y': y,
                    'font_size': font_size,
                    'font_family': font_family,
                    'element': text_elem
                })
                
            # Also check tspan elements within text
            elif text_elem.tag.endswith('}tspan') or text_elem.tag == 'tspan':
                text_content = text_elem.text or ''
                if text_content.strip():
                    # tspan might inherit position from parent text element
                    # Since ElementTree doesn't have getparent(), we'll use the element's own coordinates
                    # or default to 0 if not available
                    x = float(text_elem.get('x', 0))
                    y = float(text_elem.get('y', 0))
                    
                    text_elements.append({
                        'text': text_content.strip(),
                        'x': x,
                        'y': y,
                        'font_size': text_elem.get('font-size', ''),
                        'font_family': text_elem.get('font-family', ''),
                        'element': text_elem
                    })
    
    except Exception as e:
        logger.warning(f"Error extracting SVG text elements: {e}")
    
    return text_elements

def _find_best_svg_text_match(shape_text, svg_text_elements) -> Optional[Dict[str, Any]]:
    """Find the best matching SVG text element for a given shape text."""
    if not svg_text_elements or not shape_text.strip():
        return None
    
    shape_text_clean = shape_text.strip().lower()
    best_match = None
    best_score = 0
    
    for svg_element in svg_text_elements:
        svg_text = svg_element['text'].strip().lower()
        
        if not svg_text:
            continue
        
        # Try different matching strategies
        if fuzz:
            # Use fuzzy matching if available
            ratio = fuzz.ratio(shape_text_clean, svg_text)
            partial_ratio = fuzz.partial_ratio(shape_text_clean, svg_text)
            token_ratio = fuzz.token_set_ratio(shape_text_clean, svg_text)
            
            # Use the best of the three ratios
            score = max(ratio, partial_ratio, token_ratio)
        else:
            # Fallback to exact matching
            if shape_text_clean == svg_text:
                score = 100
            elif svg_text in shape_text_clean or shape_text_clean in svg_text:
                score = 80
            else:
                score = 0
        
        # Consider partial matches with significant overlap
        if score > best_score and score >= 70:  # Minimum threshold
            best_score = score
            best_match = {
                'svg_element': svg_element,
                'match_score': score,
                'shape_text': shape_text,
                'svg_text': svg_element['text']
            }
    
    return best_match

def _apply_coordinate_validation(shape: SlideShape, match_result, transform_info) -> SlideShape:
    """Apply coordinate validation based on SVG text matching."""
    if not match_result or not transform_info:
        shape.validation_status = "unvalidated"
        shape.validation_details = "No matching SVG text found or transform failed"
        return shape
    
    try:
        svg_element = match_result['svg_element']
        match_score = match_result['match_score']
        
        # Extract SVG coordinates
        svg_x = svg_element['x']
        svg_y = svg_element['y']
        
        # Transform SVG coordinates back to slide coordinates
        scale_x = transform_info.get('scale_x', 1.0)
        scale_y = transform_info.get('scale_y', 1.0)
        offset_x = transform_info.get('offset_x', 0)
        offset_y = transform_info.get('offset_y', 0)
        
        # Calculate expected slide coordinates
        expected_x = (svg_x - offset_x) / scale_x if scale_x != 0 else svg_x
        expected_y = (svg_y - offset_y) / scale_y if scale_y != 0 else svg_y
        
        # Compare with actual shape coordinates (convert EMU to pixels for comparison)
        EMU_PER_PIXEL = 12700  # Approximate EMU to pixel conversion
        shape_x_px = shape.x_coordinate / EMU_PER_PIXEL
        shape_y_px = shape.y_coordinate / EMU_PER_PIXEL
        
        # Calculate coordinate differences
        diff_x = abs(expected_x - shape_x_px)
        diff_y = abs(expected_y - shape_y_px)
        
        # Determine validation status based on coordinate accuracy
        tolerance = 50  # pixels
        if diff_x <= tolerance and diff_y <= tolerance and match_score >= 90:
            shape.validation_status = "validated"
            shape.validation_details = f"Coordinates validated against SVG (match: {match_score}%, diff: ±{diff_x:.1f},±{diff_y:.1f}px)"
        elif match_score >= 80:
            shape.validation_status = "partial"
            shape.validation_details = f"Text matched but coordinates differ (match: {match_score}%, diff: ±{diff_x:.1f},±{diff_y:.1f}px)"
        else:
            shape.validation_status = "questionable"
            shape.validation_details = f"Low confidence match (match: {match_score}%, diff: ±{diff_x:.1f},±{diff_y:.1f}px)"
        
        # Store additional validation metadata
        # Note: metadata field doesn't exist in SlideShape model, so we skip this for now
        # This could be added to the model later if needed
        logger.debug(f"Validation metadata: svg_match_score={match_score}, "
                    f"svg_coordinates=({svg_x}, {svg_y}), "
                    f"coordinate_difference=({diff_x}, {diff_y}), "
                    f"svg_text_matched='{svg_element['text']}')")
        
    except Exception as e:
        logger.warning(f"Error applying coordinate validation: {e}")
        shape.validation_status = "error"
        shape.validation_details = f"Validation error: {str(e)}"
    
    return shape

async def validate_coordinates_with_svg(shapes_data: List[SlideShape], svg_path: str, slide_width_emu: int, slide_height_emu: int) -> List[SlideShape]:
    """
    Validates extracted shape coordinates against the generated SVG file.
    Attempts to match text content and compare coordinate positions.
    """
    if not os.path.exists(svg_path):
        return _mark_shapes_as_unvalidated(shapes_data, "svg_file_missing")
    
    try:
        # Parse the SVG file
        tree = ET.parse(svg_path)
        root = tree.getroot()
        
        # Define SVG namespaces
        namespaces = {
            'svg': 'http://www.w3.org/2000/svg',
            'xlink': 'http://www.w3.org/1999/xlink'
        }
        
        # Extract SVG dimensions and coordinate system info
        svg_info = _extract_svg_dimensions(root)
        if not svg_info:
            return _mark_shapes_as_unvalidated(shapes_data, "svg_dimensions_extraction_failed")
        
        # Convert slide dimensions from EMU to pixels for comparison
        EMU_PER_PIXEL = 12700  # Approximate conversion factor
        slide_width_px = slide_width_emu / EMU_PER_PIXEL
        slide_height_px = slide_height_emu / EMU_PER_PIXEL
        
        # Calculate coordinate transformation parameters
        transform_info = _calculate_coordinate_transform(svg_info, slide_width_px, slide_height_px)
        if not transform_info:
            return _mark_shapes_as_unvalidated(shapes_data, "coordinate_transform_calculation_failed")
        
        # Extract all text elements from the SVG
        svg_text_elements = _extract_svg_text_elements(root, namespaces)
        if not svg_text_elements:
            return _mark_shapes_as_unvalidated(shapes_data, "no_svg_text_elements_found")
        
        logger.info(f"Found {len(svg_text_elements)} text elements in SVG for validation")
        
        # Validate each shape's coordinates
        validated_shapes = []
        for shape in shapes_data:
            if not shape.text or not shape.text.strip():
                # Skip shapes without text
                shape.validation_status = "skipped"
                shape.validation_details = "No text content to validate"
                validated_shapes.append(shape)
                continue
            
            # Find the best matching SVG text element
            match_result = _find_best_svg_text_match(shape.text, svg_text_elements)
            
            # Apply coordinate validation based on the match
            validated_shape = _apply_coordinate_validation(shape, match_result, transform_info)
            validated_shapes.append(validated_shape)
        
        # Log validation summary
        validation_counts = {}
        for shape in validated_shapes:
            status = shape.validation_status
            validation_counts[status] = validation_counts.get(status, 0) + 1
        
        logger.info(f"Coordinate validation completed: {validation_counts}")
        
        return validated_shapes
        
    except Exception as e:
        logger.warning(f"Error during coordinate validation: {str(e)}", exc_info=True)
        return _mark_shapes_as_unvalidated(shapes_data, f"validation_error: {str(e)}") 